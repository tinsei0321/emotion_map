#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "Pillow"]
# ///
"""Merge slide PNGs into a PowerPoint file with optional speaker notes."""

import argparse
import re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


EMUS_PER_INCH = 914400


def parse_slide_number(filename: str) -> int | None:
    """Extract slide number from filename like '01-slide-cover.png'."""
    match = re.match(r"(\d+)-slide-", filename)
    return int(match.group(1)) if match else None


def find_prompt_file(slide_name: str, prompts_dir: Path) -> Path | None:
    """Find matching prompt .md for a slide name."""
    # Try flat structure first: 03-prompts/01-slide-cover.md
    flat = prompts_dir / f"{slide_name}.md"
    if flat.exists():
        return flat

    # Try versioned subdirectories: 03-prompts/v6/01-slide-cover.md
    for subdir in prompts_dir.iterdir():
        if subdir.is_dir():
            candidate = subdir / f"{slide_name}.md"
            if candidate.exists():
                return candidate

    return None


def extract_notes(prompt_text: str) -> str:
    """Extract NARRATIVE GOAL and SPEAKER NOTES sections from prompt."""
    notes_parts = []

    # Extract // NARRATIVE GOAL
    ng_match = re.search(r"// NARRATIVE GOAL\s*\n(.+?)(?=\n// |\n## |\Z)", prompt_text, re.DOTALL)
    if ng_match:
        notes_parts.append(f"NARRATIVE GOAL:\n{ng_match.group(1).strip()}")

    # Extract // SPEAKER NOTES
    sn_match = re.search(r"// SPEAKER NOTES\s*\n(.+?)(?=\n// |\n## |\Z)", prompt_text, re.DOTALL)
    if sn_match:
        notes_parts.append(f"SPEAKER NOTES:\n{sn_match.group(1).strip()}")

    # If neither section found but there's KEY CONTENT, use that as fallback
    if not notes_parts:
        kc_match = re.search(r"// KEY CONTENT\s*\n(.+?)(?=\n// |\n## |\Z)", prompt_text, re.DOTALL)
        if kc_match:
            notes_parts.append(f"KEY CONTENT:\n{kc_match.group(1).strip()}")

    return "\n\n".join(notes_parts) if notes_parts else ""


def get_image_aspect_ratio(image_path: Path) -> float:
    """Return width / height ratio of an image."""
    with Image.open(image_path) as img:
        return img.width / img.height


def merge_to_pptx(
    slides_dir: str,
    output: str,
    prompts_dir: str | None = None,
    title: str | None = None,
):
    """Merge PNG slides into PowerPoint with speaker notes."""
    slides_path = Path(slides_dir)
    png_files = sorted(
        [f for f in slides_path.glob("*.png") if parse_slide_number(f.name) is not None],
        key=lambda f: parse_slide_number(f.name) or 0,
    )

    if not png_files:
        print(f"No PNG files found in {slides_dir}")
        return False

    # Auto-detect prompts directory
    if prompts_dir is None:
        project_root = slides_path.parent
        auto_prompts = project_root / "03-prompts"
        if auto_prompts.exists():
            prompts_dir = str(auto_prompts)
    prompts_path = Path(prompts_dir) if prompts_dir else None

    # Dynamic slide size from first image
    first_ratio = get_image_aspect_ratio(png_files[0])
    # Standard 16:9 is 13.333 x 7.5 inches; scale height to match ratio
    target_height = Inches(7.5)
    target_width = Inches(7.5 * first_ratio)

    prs = Presentation()
    prs.slide_width = int(target_width)
    prs.slide_height = int(target_height)

    notes_count = 0
    for png in png_files:
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(
            str(png),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Add speaker notes if prompt file found
        slide_name = png.stem  # e.g. "01-slide-cover"
        if prompts_path:
            prompt_file = find_prompt_file(slide_name, prompts_path)
            if prompt_file:
                prompt_text = prompt_file.read_text(encoding="utf-8")
                notes = extract_notes(prompt_text)
                if notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes
                    notes_count += 1

    prs.save(output)
    print(f"Created PPTX: {output}")
    print(f"Slides included: {len(png_files)}")
    if prompts_path:
        print(f"Slides with notes: {notes_count}")
    return True


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Merge slide PNGs to PPTX")
    parser.add_argument("--slides", required=True, help="Directory containing PNG slides")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument("--prompts", help="Directory containing per-slide prompt .md files (default: 03-prompts/ relative to slides parent)")
    parser.add_argument("--title", help="Presentation title")
    args = parser.parse_args()

    merge_to_pptx(args.slides, args.output, args.prompts, args.title)
